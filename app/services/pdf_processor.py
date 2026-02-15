from core.embeddings import get_embeddings
from PyPDF2 import PdfReader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from services.hashing import hash_text
from services.chunking import semantic_sections
from core.vector_store import save_vector_store
from langchain_community.vectorstores import FAISS
from docx import Document
from pptx import Presentation
import pandas as pd


import os

def extract_data(file):
    ext = os.path.splitext(file.name)[1].lower()

    if ext == ".pdf":
        return extract_pdf(file)

    elif ext == ".docx":
        return extract_docx(file)

    elif ext == ".txt":
        return extract_txt(file)

    elif ext == ".csv":
        return extract_csv(file)

    elif ext == ".pptx":
        return extract_pptx(file)

    else:
        raise ValueError(f"Unsupported file type: {ext}")

def extract_pdf(pdf):
    reader = PdfReader(pdf)
    full_text = ""
    for page in reader.pages:
        t = page.extract_text()
        if t:
            full_text += t + "\n"
    return full_text

def extract_docx(file):
    doc = Document(file)
    full_text = ""
    
    for para in doc.paragraphs:
        if para.text.strip():
            full_text += para.text + "\n"
    print(full_text)
    return full_text

def extract_txt(file):
    full_text = ""
    
    content = file.read().decode("utf-8")
    if content.strip():
        full_text += content + "\n"
    
    return full_text

def extract_csv(file):
    full_text = ""
    
    df = pd.read_csv(file)
    
    for _, row in df.iterrows():
        row_text = " ".join([str(cell) for cell in row if pd.notna(cell)])
        if row_text.strip():
            full_text += row_text + "\n"
    
    return full_text


def extract_pptx(file):
    prs = Presentation(file)
    full_text = ""
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    full_text += text + "\n"
    
    return full_text


def processFiles(files, vector_store):
    existing_chunk_hashes = set()
    existing_docs = {}

    if vector_store:
        for k, doc in vector_store.docstore._dict.items():
            ch = doc.metadata.get("chunk_hash")
            if ch:
                existing_chunk_hashes.add(ch)
                existing_docs[ch] = k

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200
    )

    all_texts = []
    all_metadatas = []
    new_chunk_hashes = set()

    for file in files:
        file_bytes = file.read()
        file_hash = hash_text(file_bytes)
        print(file_hash)
        # exit(1)
        file.seek(0)

        full_text = extract_data(file)

        sections = semantic_sections(full_text)

        for section in sections:
            chunks = (
                [section]
                if len(section) <= 1200
                else splitter.split_text(section)
            )

            for chunk in chunks:
                chunk_hash = hash_text(chunk)
                new_chunk_hashes.add(chunk_hash)

                if chunk_hash in existing_chunk_hashes:
                    continue

                all_texts.append(chunk)
                all_metadatas.append({
                    "file_hash": file_hash,
                    "chunk_hash": chunk_hash,
                    "source": file.name
                })
    # st.write(all_texts)
    # exit(1)

    # delete stale chunks
    if vector_store:
        stale_ids = [
            k for ch, k in existing_docs.items()
            if ch not in new_chunk_hashes
        ]
        print(stale_ids)

        if stale_ids:
            vector_store.delete(ids=stale_ids)

    # add new chunks
    if all_texts:
        vector_store = save_vector_store(vector_store, all_texts, all_metadatas)
        message = {"success" : "Document updated incrementally (no duplication)."}
    else:
        message = {"info" : "No new content detected."}

    return vector_store, message